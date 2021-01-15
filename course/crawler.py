import os
import re
import time
import random
import requests
import threadpool
import traceback
from multiprocessing import Pool
import pptx
import shutil
from pptx.util import Inches


class Crawler:
    def __init__(self, course_name, course_uri, retry=10, thread_number=5):
        self.retry = retry  # 失败重尝次数
        self.thread_number = thread_number  # 线程池大小
        self.course_name = course_name.strip()
        self.course_uri = course_uri.strip()
        self.base_path = "D:\\video1"
        self.site = "http://cache.gensee.com"
        self.video_path = os.path.join(self.base_path, self.course_name)
        self.ts_path = os.path.join(self.video_path, 'ts')  # 视频片段路径
        if not os.path.exists(self.ts_path):
            os.makedirs(self.ts_path)

        self.timeout = 10
        self.session = requests.Session()
        self.session.headers = {
            "Connection": "close",
            "Host": "cache.gensee.com",
            "Origin": "http://edusoho.gensee.com",
            "User-Agent": "Mozilla/5.0 (Windows NT 10.0; Win64; x64; rv:84.0) Gecko/20100101 Firefox/84.0"
        }

    # 获取课程相关的所有ts分片视频
    def generate_ts_url(self):
        resp = self.session.get(self.site + '/' + self.course_uri + '/hls/record.json').json()
        if not resp['m3u8']:
            raise Exception("Not found m3u8 resource from this site.")
        for item in resp['m3u8']:
            if 'file' in item:
                resp = self.session.get(self.site + '/' + self.course_uri + '/hls/' + item['file'])
                for line in str.split(resp.text, '\n'):
                    if re.match(r".*ts", line):
                        yield self.site + '/' + self.course_uri + '/hls/' + line.strip()

    # 下载ts 视频，失败重试，且最多重试self.retry次
    def download_ts(self, url):
        retry = 0
        ts_file = os.path.basename(url)
        while True:
            try:
                print(f"【download_ts】start download. ts_file: {ts_file}, retry time: {retry + 1}, ts_url: {url}")
                if retry >= self.retry:
                    print(f"【download_ts】download failed. error: retry failed retry, ts_file: {ts_file}")
                    break

                time.sleep(random.random())
                data = self.session.get(url, stream=True, timeout=self.timeout)
                with open(os.path.join(self.ts_path, ts_file), 'wb') as f:
                    f.write(data.content)
                    print(f"【download_ts】download success. ts_file: {ts_file}. ")
                return
            except Exception as err:
                print(f"【download_ts】download exception. ts_file: {ts_file}, retry time: {retry + 1}, err: {err}")
                retry += 1

    # 合并ts视频到video.mp4
    def merge_ts(self):
        def sort_ts(x):
            a, b = re.findall(r'\d+', x)
            return int(a), int(b)

        ts_files = sorted(os.listdir(self.ts_path), key=sort_ts)
        for ts_file in ts_files:
            print(f"【merge_ts】merge ts_file: {ts_file}")
            with open(os.path.join(self.ts_path, ts_file), 'rb') as f:
                content = f.read()
            with open(os.path.join(self.video_path, 'video.mp4'), 'ab') as v:
                v.write(content)
        if os.path.exists(self.ts_path):
            shutil.rmtree(self.ts_path)

    def download_pic(self):
        pic_path_list = []
        resp = requests.get(self.site + '/' + self.course_uri + '/record.xml.js').json()
        for module in resp['conf']['module']:
            if module['name'] == 'document':
                documents = module['document']
                if type(documents) is not list:
                    documents = [documents]
                for document in documents:
                    pages = document['page']
                    ppt_name = document['id'] + '-' + os.path.splitext(document['name'])[0]
                    pic_path = os.path.join(self.video_path, ppt_name)
                    if not os.path.exists(pic_path):
                        os.makedirs(pic_path)
                    if type(pages) is not list:
                        pages = [pages]
                    for page in pages:
                        # title = page.get('title')
                        # title = str(int(page.get('id')) + 1) + '-' + title.replace('&', '和')
                        title = str(int(page.get('id')) + 1)
                        print(f"【download_pic】start download. pic_path: {pic_path}, title: {title}")
                        resp = self.session.get(self.site + '/' + self.course_uri + '/' + page.get('hls'))
                        with open(os.path.join(pic_path, title + ".png"), 'wb') as f:
                            f.write(resp.content)
                    pic_path_list.append({'name': ppt_name, 'path': pic_path})
        return pic_path_list

    def merge_pic(self, pic_path_list):
        def sort_pic(x):
            r = re.findall(r'\d+', x)
            return int(r[0])

        for pic_path in pic_path_list:
            ppt_file = pptx.Presentation()
            pic_files = sorted(os.listdir(pic_path['path']), key=sort_pic)

            for pic_file in pic_files:
                print(f"【merge_pic】merge pic_file: {pic_file}")
                pic = os.path.join(pic_path['path'], pic_file)
                slide = ppt_file.slides.add_slide(ppt_file.slide_layouts[1])
                # slide.shapes.placeholders[0].text = os.path.splitext(os.path.basename(pic))[0]
                slide.shapes.placeholders[0].text = ""
                slide.shapes.add_picture(pic, Inches(0), Inches(0), Inches(10), Inches(7.5))
            ppt_file.save(os.path.join(self.video_path, f"{pic_path['name']}.pptx"))
            if os.path.exists(pic_path['path']):
                shutil.rmtree(pic_path['path'])

    def run(self):
        try:
            # print(f"【下载TS文件】课程名称：{self.course_name}, 课程uri: {self.course_uri}, 开始下载 ...")
            # executor = threadpool.ThreadPool(self.thread_number)
            # tasks = threadpool.makeRequests(self.download_ts, self.generate_ts_url())
            # [executor.putRequest(task) for task in tasks]
            # executor.wait()
            # print(f"【下载TS文件】成功下载TS文件. TS文件路径: {self.ts_path}.")
            #
            # print(f"【合并ts文件】TS文件路径: {self.ts_path}，开始合并 ...")
            # self.merge_ts()
            # print(f"【合并ts文件】成功合并TS文件. TS文件路径: {self.ts_path}")

            print(f"【下载课件】课程名称：{self.course_name}, 课程uri: {self.course_uri}, 开始下载 ...")
            pic_path_list = self.download_pic()
            print(f"【下载课件】成功下载课件. ")

            print(f"【合并课件】课件路径: {pic_path_list}，开始合并 ...")
            self.merge_pic(pic_path_list)
            print(f"【合并课件】成功合并课件. 课件路径: {pic_path_list}")
            if os.path.exists(self.ts_path):
                shutil.rmtree(self.ts_path)
        except Exception as err:
            print(f"【下载失败】课程名称：{self.course_name}, 课程uri: {self.course_uri}，err: {err}")
            traceback.print_exc()


if __name__ == '__main__':
    course_list = [
        {"name": "任务6-1：沙盘游戏咨询的个案督导（尹立）6月5日",
         "uri": "gsgetrecord/recordcz132.gensee.net/gsrecord/33162/sbr/2020_06_05/4ovLqTEFrr_1591355032"},
    ]
    # for script in course_list:
    #     Crawler(script['name'], script['uri']).run()
    #     break
    pool = Pool(5)
    for course in course_list:
        pool.apply_async(Crawler(course['name'], course['uri']).run, args=())
        print('Task %s has been submited' % course['name'])
        # break
    print('Waiting for all subprocesses done...')
    pool.close()
    pool.join()
    print('All subprocesses done.')
